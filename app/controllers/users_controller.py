import os
import re
import fitz  
import openpyxl  
from pptx import Presentation
from io import BytesIO
from flask import current_app, request, jsonify,render_template,redirect,url_for
from werkzeug.utils import secure_filename
from docx import Document
from config import Config 
from app.models.category_models import Category
from app.models.file_models import File
from flask_login import current_user
from app.models.user_models import User
from app.services.services import CRUDService
from app.utils.validation_utils import Validation
ALLOWED_EXTENSIONS = {"txt", "csv", "pdf", "docx", "xlsx", "pptx"}


file_service = CRUDService(File)
user_services = CRUDService(User)
categories_service = CRUDService(Category)

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_xlsx(file):
    """Extract text from an XLSX file."""
    text = []
    file.seek(0)  
    workbook = openpyxl.load_workbook(file, data_only=True)
    
    for sheet in workbook.sheetnames:
        ws = workbook[sheet]
        for row in ws.iter_rows(values_only=True):
            text.append(" ".join(str(cell) for cell in row if cell)) 
            
    return "\n".join(text).strip()

def extract_text_from_docx(file):
    """Extract text from a DOCX file."""
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pdf(file):
    """Extract text from a PDF file."""
    text = ""
    pdf_bytes = file.read()
    pdf_stream = BytesIO(pdf_bytes)
    
    try:
        with fitz.open(stream=pdf_stream, filetype="pdf") as doc:
            for page in doc:
                text += page.get_text("text") + "\n"
    except Exception as e:
        print("Error reading PDF:", e)
        return ""

    return text.strip()

def extract_text_from_pptx(file):
    """Extract text from a PPTX file."""
    text = []
    ppt = Presentation(file)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text).strip()

def get_categories_from_db():
    categories = {}
    db_categories = categories_service.get()

    for category in db_categories:
        if category.keywords:
            categories[category.name] = category.keywords  

    return categories

def categorize_content(content):
    if not content.strip():  
        return "Uncategorized"

    categories = get_categories_from_db() 

    for category, keywords in categories.items():
        for keyword in keywords:
            if re.search(rf"\b{keyword}\b", content, re.IGNORECASE):
                return category  

    return "Uncategorized"

def upload_file(request):
    uploader_id = current_user.id if current_user.is_authenticated else None

    if "docs_file" not in request.files:
        return jsonify({'error': True, 'message': 'No file part'}), 400

    file = request.files["docs_file"]
    
    if file.filename == "" or not allowed_file(file.filename):
        return jsonify({'error': True, 'message': 'Invalid file'}), 400

   # Extract filename and extension
    filename_with_ext = secure_filename(file.filename)
    filename, file_ext = os.path.splitext(filename_with_ext) 
    file_ext = file_ext.lower().lstrip(".")  

    content = ""
    if file_ext == "pdf":
        content = extract_text_from_pdf(file)
    elif file_ext == "docx":
        content = extract_text_from_docx(file)
    elif file_ext == "xlsx":
        content = extract_text_from_xlsx(file)
    elif file_ext == "pptx":
        content = extract_text_from_pptx(file) 
    else:
        content = file.read().decode("utf-8")


    # Categorize content
    category_name = categorize_content(content)
    if category_name == "Uncategorized":
        return jsonify({'error': True, 'message': 'File does not match any category. Upload denied.'}), 400

    # Fetch category from database
    category = categories_service.get_one(name=category_name)
    if not category:
        return jsonify({'error': True, 'message': 'Category not found in database.'}), 400

    
    existing_file = file_service.get_one(
        filename=filename, category_id=category.id, uploader_id=uploader_id
    )
    if existing_file:
        return jsonify({
            'error': True,
            'message': f'You have already uploaded a file named "{filename}" under the "{category_name}" category!'
        }), 400

    upload_folder = os.path.join(current_app.root_path, "static", "upload")
    os.makedirs(upload_folder, exist_ok=True)  

    filepath = os.path.join("static", "upload", filename_with_ext)
    full_path = os.path.join(upload_folder, filename_with_ext)

    file.seek(0, os.SEEK_END)
    file_size = file.tell()
    file.seek(0)
    file.save(full_path)

    uploader_id = current_user.id if current_user.is_authenticated else None
    
    file_data = {
        "filename": filename,
        "filepath": filepath,
        "file_size": file_size,
        "file_extension": file_ext,
        "category_id": category.id,
        "uploader_id": uploader_id 
    }
    file_service.create(**file_data)

    return jsonify({
        "filename": filename,
        "category": category_name,
        "file_url": f"/{filepath}",
        "message": "File uploaded successfully!"
    })


def delete_file(request, file_id=None):
    if request.method == 'DELETE':
        file = file_service.get_one(id=file_id)
        
        if not file:
            return jsonify({"success": False, "message": "File not found"}), 404

        # Remove file from storage
        try:
            file_path = os.path.join(current_app.root_path, file.filepath)
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception as e:
            return jsonify({"success": False, "message": f"Error deleting file: {str(e)}"}), 500

        file_service.delete(file.id)

        return jsonify({"success": True, "message": "File deleted successfully!"})
    
def docs_list(request):
    if request.method == 'GET':
        categories = categories_service.get()
        return render_template('users/docs.html', categories=categories)
    

def dashboard_report_users():
    if not current_user.is_authenticated:
        return redirect(url_for('auth.login'))
    
    # Fetch categories
    categories = categories_service.get()
    files = file_service.get(uploader_id=current_user.id)  

    
    file_data = {category.name: 0 for category in categories}  

    for file in files:
        if file.category_id in [category.id for category in categories]:
            category_name = next(cat.name for cat in categories if cat.id == file.category_id)
            file_data[category_name] += 1
            
    print(file_data)
    return render_template(
        'users/dashboard.html',
        file_data=file_data  
    )
    


def update_profiles():
    data = request.json
    user_id = current_user.id
    user = user_services.get_one(id=user_id)

    if not user:
        return jsonify({"success": False, "message": "User not found"}), 404

    # Fields to update
    update_data = {
        "firstname": data.get("asfname"),
        "middlename": None if data.get("asmname") == "N/A" else data.get("asmname"),
        "lastname": data.get("aslname"),
    }

    # Handle password update
    current_password = data.get("currentPassword")
    new_password = data.get("newPassword")

    if new_password:  # User wants to change password
        if not current_password:
            return jsonify({"success": False, "message": "Current password is required to change password"}), 400
        
        if not user.check_password(current_password):  # Use model method
            return jsonify({"success": False, "message": "Incorrect current password"}), 400

         # Validate new password
        password_error = Validation.is_strong_password(new_password)
        if password_error:
            return jsonify({"success": False, "message": password_error}), 400

        user.set_password(new_password)  # Use model method

    # Update user information
    updated_user = user_services.update(user_id, **update_data)

    if updated_user:
        return jsonify({"success": True, "message": "Profile updated successfully"})
    else:
        return jsonify({"success": False, "message": "Error updating profile"}), 500
   