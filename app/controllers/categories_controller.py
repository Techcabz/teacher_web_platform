from flask import request, render_template, jsonify
from app.services.services import CRUDService
from app.models.category_models import Category
from app.utils.validation_utils import Validation
import os
import re
import fitz  
import openpyxl  
import xlrd
from pptx import Presentation
from io import BytesIO
from docx import Document
from werkzeug.utils import secure_filename

ALLOWED_EXTENSIONS = {"txt", "csv", "pdf", "docx", "xlsx", "pptx", "xls"}

category_service = CRUDService(Category)
def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_xlsx(file):
    """Extract text from an XLSX file."""
    text = []
    file.seek(0)  
    workbook = openpyxl.load_workbook(file, data_only=True)
    
    for sheet in workbook.sheetnames:
        ws = workbook[sheet]
        for row in ws.iter_rows(values_only=True):
            text.append(" ".join(str(cell) for cell in row if cell)) 
            
    return "\n".join(text).strip()

def extract_text_from_xls(file):
    """Extract text from an XLS file."""
    text = []
    file.seek(0)
    workbook = xlrd.open_workbook(file_contents=file.read())
    
    for sheet in workbook.sheets():
        for row_idx in range(sheet.nrows):
            row_values = sheet.row_values(row_idx)
            text.append(" ".join(str(cell) for cell in row_values if cell))

    return "\n".join(text).strip()

def extract_text_from_docx(file):
    """Extract text from a DOCX file."""
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pdf(file):
    """Extract text from a PDF file."""
    text = ""
    pdf_bytes = file.read()
    pdf_stream = BytesIO(pdf_bytes)
    
    try:
        with fitz.open(stream=pdf_stream, filetype="pdf") as doc:
            for page in doc:
                text += page.get_text("text") + "\n"
    except Exception as e:
        print("Error reading PDF:", e)
        return ""

    return text.strip()

def extract_text_from_pptx(file):
    """Extract text from a PPTX file."""
    text = []
    ppt = Presentation(file)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text).strip()
from flask import request, render_template, jsonify
from werkzeug.utils import secure_filename
import os

def categories_keyword(request):
    if request.method == 'GET':
        return render_template('admin/keyword.html', content=None)

    elif request.method == 'POST':
        if "docs_file" not in request.files:
            return jsonify({'error': True, 'message': 'No file part'}), 400

        file = request.files["docs_file"]

        if file.filename == "" or not allowed_file(file.filename):
            return jsonify({'error': True, 'message': 'Invalid file'}), 400

        filename_with_ext = secure_filename(file.filename)
        filename, file_ext = os.path.splitext(filename_with_ext)
        file_ext = file_ext.lower().lstrip(".")

        content = ""
        if file_ext == "pdf":
            content = extract_text_from_pdf(file)
        elif file_ext == "docx":
            content = extract_text_from_docx(file)
        elif file_ext == "xlsx":
            content = extract_text_from_xlsx(file)
        elif file_ext == "xls":
            content = extract_text_from_xls(file)
        elif file_ext == "pptx":
            content = extract_text_from_pptx(file)
        else:
            content = file.read().decode("utf-8", errors="ignore")

        return render_template('admin/keyword.html', content=content,filename=filename_with_ext)

    

         
def categories(request, category_id=None):
    if request.method == 'GET':
        if category_id:
            category = category_service.get_one(id=category_id)
            if not category:
                return jsonify({'success': False, 'message': 'Category not found'}), 404
            return jsonify(category.serialize()), 200
        
        categories = category_service.get()
        return render_template('admin/category.html', categories=categories)

    elif request.method == 'POST':
        name = request.form['ciname']
        keywords = request.form.getlist('keywords[]')

        if not keywords or all(k.strip() == "" for k in keywords):
            return jsonify({'success': False, 'message': 'Keywords are required.'}), 400


        # Validate category name
        if Validation.has_repeated_characters(name):
            return jsonify({'success': False, 'message': 'Name must not have three or more consecutive repeated characters.'}), 400
        
        existing_category = category_service.get_one(name=name)
        if existing_category:
            return jsonify({'success': False, 'message': 'Category name is already taken.'}), 400

        # Ensure no duplicate keywords within the same submission
        if len(set(keywords)) != len(keywords):
            return jsonify({'success': False, 'message': 'Duplicate keywords found in input.'}), 400

        new_category = category_service.create(name=name, keywords=keywords)
        if new_category:
            return jsonify({'success': True, 'message': 'Category created successfully!'}), 201

    elif request.method == 'PUT':
        if category_id:
            name = request.form['ciname']
            keywords = request.form.getlist('keywords[]')  # Keywords from request

            existing_category = category_service.get_one(id=category_id)
            
            if not existing_category:
                return jsonify({'success': False, 'message': 'Category not found.'}), 404

            # Extract existing keywords
            existing_keywords = set(existing_category.keywords if existing_category.keywords else [])

            # Remove duplicates in input
            unique_keywords = set(keywords)  

            # Compare old and new keywords
            if existing_keywords == unique_keywords:
                return jsonify({'success': False, 'message': 'No changes detected. The provided keywords are the same.'}), 400

            # ✅ Update category with the new set of keywords
            updated_category = category_service.update(category_id, name=name, keywords=list(unique_keywords))
            
            if updated_category:
                return jsonify({'success': True, 'message': 'Category updated successfully!'}), 200

        return jsonify({'success': False, 'message': 'Error updating category.'}), 500

    elif request.method == 'DELETE':
        if category_id:
            existing_category = category_service.get_one(id=category_id)
            if not existing_category:
                return jsonify({'success': False, 'message': 'Category not found.'}), 404
            delete_success = category_service.delete(category_id)
            if delete_success:
                return jsonify({'success': True, 'message': 'Category deleted successfully!'}), 200
            return jsonify({'success': False, 'message': 'Error deleting category.'}), 500

    return jsonify({'success': False, 'message': 'Create method must be POST.'}), 405
